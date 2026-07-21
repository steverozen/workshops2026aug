#!/usr/bin/env python3
"""Build the Wednesday-morning trajectory-inference slides.

Format is copied from
slides/tues_afternoon/Steve_Rozen_R_CITEseq_and_scATACseq_2026_08_04.pptx:
16:9, Arial, plain black text boxes on white, 22pt body, 28pt section
dividers, no bullet glyphs, blank paragraph between items.
"""

import os

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = ("/home/steve/github/workshops2026aug/slides/wed_morning/"
       "Steve_Rozen_R_trajectory_inference_2026_08_05.pptx")

# Geometry lifted from the reference deck (EMU).
SLIDE_W, SLIDE_H = 12192000, 6858000
BODY_X, BODY_Y = 961560, 1272600
BODY_W, BODY_H = 9558000, 4112280
DIV_X, DIV_Y = 3637080, 1771200
DIV_W, DIV_H = 10027800, 3434760

FONT = "Arial"
BLACK = RGBColor(0, 0, 0)

prs = Presentation()
prs.slide_width = Emu(SLIDE_W)
prs.slide_height = Emu(SLIDE_H)
BLANK = prs.slide_layouts[6]


def _style(run, size, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = BLACK


def textbox(slide, x, y, w, h, lines, size=22, line_spacing=1.0):
    """lines: list of (text, bold). Empty text yields a blank paragraph."""
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        if text:
            _style(p.add_run(), size, bold)
            p.runs[0].text = text
        else:
            # keep blank paragraphs at the body size so spacing matches
            _style(p.add_run(), size)
            p.runs[0].text = ""
    return box


def divider(text):
    s = prs.slides.add_slide(BLANK)
    textbox(s, DIV_X, DIV_Y, DIV_W, DIV_H, [(text, False)], size=28)
    return s


def content(heading, items, size=22):
    """heading then a blank line then one blank-separated line per item."""
    s = prs.slides.add_slide(BLANK)
    lines = [(heading, False), ("", False)]
    for it in items:
        lines.append((it, False))
        lines.append(("", False))
    textbox(s, BODY_X, BODY_Y, BODY_W, BODY_H, lines[:-1], size=size)
    return s


def three_col(heading, rows, size=16):
    """Aligned three-column comparison built from three plain text boxes.

    Every cell is kept short enough to stay on one line so the columns
    line up without needing a table.
    """
    s = prs.slides.add_slide(BLANK)
    textbox(s, BODY_X, 640000, BODY_W, 500000, [(heading, False)], size=22)
    widths = [2750000, 3900000, 3900000]
    xs = [BODY_X]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)
    for col in range(3):
        lines = [(r[col], r[3]) for r in rows]
        textbox(s, xs[col], BODY_Y + 200000, widths[col], BODY_H,
                lines, size=size, line_spacing=1.5)
    return s


FIG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "figures")
FIG_TOP = 1272600
FIG_MAX_W = 11000000
FIG_MAX_H = 4900000


def figure(heading, png, caption=None):
    """Heading line, then the figure scaled to fit and centred."""
    s = prs.slides.add_slide(BLANK)
    textbox(s, BODY_X, 640000, BODY_W, 500000, [(heading, False)], size=22)
    pic = s.shapes.add_picture(os.path.join(FIG_DIR, png), 0, 0)
    # add_picture uses the PNG's native size, so rescale it to the box
    scale = min(FIG_MAX_W / pic.width, FIG_MAX_H / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int((SLIDE_W - pic.width) / 2)
    pic.top = FIG_TOP
    if caption:
        textbox(s, BODY_X, FIG_TOP + pic.height + 90000, BODY_W, 400000,
                [(caption, False)], size=14)
    return s


SCHEMATIC = "Schematic, drawn to show the algorithm. Not package output."


# ---------------------------------------------------------------- slides

divider("Trajectory inference: Slingshot and Monocle 3")

content("What we are trying to do:", [
    "Cells in a differentiation process are captured at all stages at once",
    "We have no time axis, only a snapshot",
    "Trajectory inference orders cells along a path of gradual change",
    "Pseudotime = how far along that path a cell sits",
    "It is a distance, not a clock. Units are arbitrary.",
])

content("The two tools, in one sentence each:", [
    "Slingshot fits smooth curves through a cluster-level tree, and gives "
    "one pseudotime per lineage",
    "Monocle 3 fits a principal graph to the UMAP embedding, and gives one "
    "pseudotime per cell, measured from a root you pick",
])

content("Slingshot: how it works", [
    "Input: a reduced-dimension embedding plus a vector of cluster labels. "
    "That is all.",
    "Stage 1, getLineages(): minimum spanning tree on the cluster centroids, "
    "which gives the branching structure",
    "Stage 2, getCurves(): fit simultaneous principal curves, one per lineage",
    "The curves are shrunk toward each other near branch points, so lineages "
    "share a trunk and then separate",
])

figure("Slingshot stage 1: the cluster-level tree",
       "fig_cluster_tree.png", SCHEMATIC)

figure("Slingshot stage 2: one smooth curve per lineage",
       "fig_slingshot_curves.png",
       "The two curves nearly coincide along the trunk, then separate. "
       + SCHEMATIC)

content("Monocle 3: how it works", [
    "Input: a cell_data_set, built through Monocle's own preprocessing",
    "cluster_cells() assigns both clusters and partitions",
    "learn_graph() fits a principal graph by reversed graph embedding, "
    "within each partition",
    "order_cells() sets the root, and pseudotime is the distance along the "
    "graph from that root",
])

figure("Monocle 3: the principal graph, fitted to the embedding",
       "fig_monocle_graph.png",
       "Grey spokes show each cell assigned to its nearest vertex. Pseudotime "
       "is distance from the root along the black edges. " + SCHEMATIC)

content("Careful: the black dots mean different things", [
    "Slingshot's dots are cluster centroids. Eight clusters, eight dots. "
    "You can label each one with a cell type.",
    "Monocle 3's dots are principal points: landmarks the algorithm invents "
    "to build a skeleton through the cells",
    "They are not cells and not clusters, and no single one means anything "
    "biologically. Only the shape of the whole graph does.",
    "How many is set by ncenter in learn_graph_control, usually far more "
    "than you have clusters",
    "They are what you name in order_cells(root_pr_nodes = ...), hence "
    "labels like Y_1 and Y_17",
], size=20)

figure("The same picture, two different meanings",
       "fig_vertices.png",
       "Left: one dot per cell type. Right: 28 landmarks, chosen by ncenter, "
       "individually meaningless. " + SCHEMATIC)

figure("The same data, the two models",
       "fig_contrast.png", SCHEMATIC)

three_col("Side by side:", [
    # non-breaking space, an empty cell here collapses and the column
    # would then sit one line higher than the other two
    (" ", "Slingshot", "Monocle 3", True),
    ("Topology", "MST on clusters, then curves", "Principal graph (SimplePPT)", False),
    ("Fitted in", "Any reduced space, e.g. 10 PCs", "The UMAP embedding, usually 2D", False),
    ("Pseudotime", "Matrix: cells x lineages", "One vector per cell", False),
    ("Measured from", "A start cluster you supply", "Root nodes you supply", False),
    ("Branches", "Explicit, named lineages", "Structure in the graph", False),
    ("Disconnected parts", "No, it connects everything", "Yes, via partitions", False),
    ("Cycles", "No, tree only", "Yes, loops can be closed", False),
    ("DE on pseudotime", "Not included, use tradeSeq", "Built in: graph_test()", False),
    ("Ecosystem", "One Bioconductor package", "Framework, own object type", False),
])

content("Difference 1: curves versus a graph", [
    "Slingshot gives a few smooth 1-D paths. That is the developmental "
    "picture most people already have in mind.",
    "Monocle 3 gives a piecewise-linear skeleton that can branch repeatedly, "
    "close loops, and be disconnected",
    "More flexible topology, but less smooth and harder to summarise",
])

content("Difference 2: which space the fit happens in", [
    "This is the one with the biggest practical consequences",
    "Slingshot can fit in, say, 10 PCs and only project to 2D for the figure",
    "Monocle 3 fits the graph in the UMAP embedding itself",
    "UMAP does not preserve global distances, so the trajectory inherits "
    "whatever distortion UMAP introduced",
    "Change the UMAP seed or n_neighbors and the trajectory can change shape",
])

content("Difference 3: per-lineage versus single pseudotime", [
    "Slingshot: a cell before the branch point gets a pseudotime in every "
    "downstream lineage, plus a weight for each",
    "That is exactly what tradeSeq consumes to ask whether a gene behaves "
    "differently in lineage 1 than in lineage 2",
    "Monocle 3: one number per cell, so branch-specific questions mean "
    "subsetting and re-running",
])

content("Difference 4: composable versus all-in-one", [
    "Slingshot needs only an embedding and cluster labels, so it drops into "
    "a Seurat or scran workflow you already trust",
    "Monocle 3 wants you inside its cell_data_set world: its normalization, "
    "its clustering, its batch alignment",
    "Convenient if you start there, awkward if you are converting from Seurat",
])

content("Gotchas to watch for:", [
    "Slingshot will connect anything you give it. Feed it T cells and "
    "monocytes together and it will draw a lineage between them.",
    "So subset to related cells first, and set start.clus explicitly",
    "Monocle 3 infers trajectories within a partition. Cells in other "
    "partitions get infinite pseudotime, which is the usual cause of a "
    "half-grey UMAP.",
    "Neither tool finds the root for you. Direction comes from biology.",
])

content("Which to use:", [
    "Start with Slingshot: the two stages are easy to reason about, and it "
    "composes with the Seurat workflow we used earlier",
    "Reach for Monocle 3 when the topology is genuinely complex, or when you "
    "want the built-in DE and gene modules",
    "In the dynverse benchmark Slingshot was among the top performers for "
    "linear and branching topologies",
    "Note that the benchmark tested Monocle 2 (DDRTree), not Monocle 3, so "
    "that result does not transfer directly",
])

content("Hands-on tutorials:", [
    "Slingshot, NBIS workshop-scRNAseq, Seurat-based, bone marrow data, "
    "Slingshot plus tradeSeq:",
    "https://nbisweden.github.io/workshop-scRNAseq/labs/seurat/"
    "seurat_07_trajectory.html",
    "Monocle 3, Galaxy Training Network, about 3 hours in R:",
    "https://training.galaxyproject.org/training-material/topics/single-cell/"
    "tutorials/scrna-case_monocle3-rstudio/tutorial.html",
], size=18)

content("References:", [
    "Street et al. 2018, Slingshot: cell lineage and pseudotime inference "
    "for single-cell transcriptomics, BMC Genomics. "
    "https://doi.org/10.1186/s12864-018-4772-0",
    "Cao et al. 2019, The single-cell transcriptional landscape of mammalian "
    "organogenesis, Nature. https://doi.org/10.1038/s41586-019-0969-x",
    "Saelens et al. 2019, A comparison of single-cell trajectory inference "
    "methods, Nature Biotechnology. https://doi.org/10.1038/s41587-019-0071-9",
    "Van den Berge et al. 2020, Trajectory-based differential expression "
    "analysis for single-cell sequencing data, Nature Communications. "
    "https://doi.org/10.1038/s41467-020-14766-3",
], size=16)

prs.save(OUT)
print("wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
