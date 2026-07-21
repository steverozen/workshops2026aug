#!/usr/bin/env python3
"""Build the opening single-cell-bioinformatics overview slides.

Outline-only first pass: a map of the field along two axes, what the
technology measures and what question you are asking of it, so each
hands-on tutorial can be placed at an intersection.

Format is copied from
slides/wed_morning/make_traj_slides.py, which in turn took it from the
Tuesday-afternoon CITE-seq deck: 16:9, Arial, plain black text boxes on
white, 22pt body, 28pt section dividers, no bullet glyphs, blank
paragraph between items.
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = ("/home/steve/github/workshops2026aug/slides/tues_morning/"
       "single-cell-bioinformatics-overview.pptx")

# Geometry lifted from the reference deck (EMU).
SLIDE_W, SLIDE_H = 12192000, 6858000
BODY_X, BODY_Y = 961560, 1272600
BODY_W, BODY_H = 9558000, 4112280
DIV_X, DIV_Y = 3637080, 1771200
DIV_W, DIV_H = 10027800, 3434760

FONT = "Arial"
BLACK = RGBColor(0, 0, 0)

prs = Presentation()
prs.slide_width = Emu(SLIDE_W)
prs.slide_height = Emu(SLIDE_H)
BLANK = prs.slide_layouts[6]


def _style(run, size, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = BLACK


def textbox(slide, x, y, w, h, lines, size=22, line_spacing=1.0):
    """lines: list of (text, bold). Empty text yields a blank paragraph."""
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        if text:
            _style(p.add_run(), size, bold)
            p.runs[0].text = text
        else:
            # keep blank paragraphs at the body size so spacing matches
            _style(p.add_run(), size)
            p.runs[0].text = ""
    return box


def divider(text):
    s = prs.slides.add_slide(BLANK)
    textbox(s, DIV_X, DIV_Y, DIV_W, DIV_H, [(text, False)], size=28)
    return s


def content(heading, items, size=22):
    """heading then a blank line then one blank-separated line per item."""
    s = prs.slides.add_slide(BLANK)
    lines = [(heading, False), ("", False)]
    for it in items:
        lines.append((it, False))
        lines.append(("", False))
    textbox(s, BODY_X, BODY_Y, BODY_W, BODY_H, lines[:-1], size=size)
    return s


def three_col(heading, rows, size=16):
    """Aligned three-column comparison built from three plain text boxes.

    Every cell is kept short enough to stay on one line so the columns
    line up without needing a table.
    """
    s = prs.slides.add_slide(BLANK)
    textbox(s, BODY_X, 640000, BODY_W, 500000, [(heading, False)], size=22)
    widths = [2750000, 3900000, 3900000]
    xs = [BODY_X]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)
    for col in range(3):
        lines = [(r[col], r[3]) for r in rows]
        textbox(s, xs[col], BODY_Y + 200000, widths[col], BODY_H,
                lines, size=size, line_spacing=1.5)
    return s


# ---------------------------------------------------------------- slides

divider("Single-cell bioinformatics, an overview")

content("How to read this workshop:", [
    "Two questions organise everything we will do",
    "First: what does the technology actually measure?",
    "Second: what are you asking of the data once you have it?",
    "Those axes are independent. The same question gets asked of very "
    "different assays, and the same assay gets asked very different questions.",
    "Every tutorial in the next three days sits at one intersection",
])

# ------------------------------------------------------------ technologies

divider("Technologies")

content("Non-spatial: cells are dissociated out of the tissue", [
    "Position is destroyed, throughput is high, and every cell is measured "
    "the same way",
    "Transcriptome: scRNA-seq. The default, and most of what follows assumes it.",
    "Chromatin accessibility: scATAC-seq",
    "Surface protein: CITE-seq, antibody-derived tags read out as extra "
    "features alongside RNA",
    "Immune receptor repertoire: single-cell TCR-seq and BCR-seq",
    "Joint multi-omic assays measure two of these in the same cell",
], size=20)

content("Cells or nuclei?", [
    "A choice you make at the bench that you then have to live with in the "
    "analysis",
    "Whole cells: more RNA per cell, more cytoplasmic and mitochondrial reads, "
    "and dissociation stress shows up as a gene programme",
    "Nuclei: works on frozen and hard-to-dissociate tissue such as brain, less "
    "RNA per nucleus, enriched for unspliced transcripts",
    "Do not integrate the two casually. The difference is larger than most "
    "batch effects.",
], size=20)

content("Perturbation screens: don't just observe, intervene", [
    "Everything else on this list is observational. You take a snapshot and "
    "infer what is driving it.",
    "Perturb-seq puts a pooled CRISPR screen and scRNA-seq in the same "
    "experiment: each cell carries a guide RNA, and the guide is read out "
    "alongside the transcriptome",
    "So the readout is not survival or one reporter, it is the whole "
    "transcriptional consequence of knocking out that gene, per cell",
    "Replogle et al. took this genome-scale: every expressed gene knocked "
    "down, across 2.5 million cells",
    "Related idea, heritable barcodes instead of guides, and you recover "
    "lineage rather than perturbation",
], size=20)

content("Why Perturb-seq changes the analysis:", [
    "The cell label is now assigned by you, not inferred. You know which gene "
    "was hit before you cluster anything.",
    "That turns 'what is regulating this?' from an inference into a "
    "measurement, and gives the gene-network methods something to be tested "
    "against",
    "New problems come with it: guide assignment is noisy, multiplets carry "
    "two guides, and knockdown is partial and variable between cells",
    "Not covered in this workshop's tutorials. Worth knowing it exists if you "
    "run screens.",
], size=20)

content("Spatial: position is kept", [
    "Molecular measurement is tied back to tissue architecture, at the cost of "
    "either gene coverage or resolution",
    "The field splits cleanly in two (Moses and Pachter 2022):",
    "Imaging-based: molecules are counted in place, one fluorescent round at a "
    "time. MERFISH / Vizgen MERSCOPE, Nanostring CosMx, 10x Xenium, and Akoya "
    "CODEX for protein.",
    "Sequencing-based: transcripts are captured onto a barcoded surface, then "
    "sequenced off the tissue. Visium, Visium HD, Slide-seq.",
], size=20)

three_col("Imaging-based versus sequencing-based:", [
    # non-breaking space, an empty cell here collapses and the column
    # would then sit one line higher than the other two
    (" ", "Imaging-based", "Sequencing-based", True),
    ("Genes measured", "Targeted panel, 100s to 5000", "Whole transcriptome", False),
    ("Resolution", "Sub-cellular, single molecule", "Spot, bead or bin", False),
    ("Cell boundaries", "Segmentation required", "Deconvolution required", False),
    ("Sensitivity", "High per gene", "Lower, capture-limited", False),
    ("Sample format", "FFPE or fresh frozen", "Mostly fresh frozen, FFPE now", False),
    ("Field of view", "Small, imaging time bound", "Whole section", False),
    ("Data you get", "Molecule coordinates", "Counts per spot, plus image", False),
    ("Main headache", "Bad segmentation", "Spots are cell mixtures", False),
    ("Examples", "MERSCOPE, CosMx, Xenium", "Visium, Visium HD, Slide-seq", False),
])

content("Where the workshop lands on this map:", [
    "Tuesday: dissociated scRNA-seq end to end, then CITE-seq and scATAC-seq",
    "Wednesday morning: differential expression, reference mapping and "
    "trajectory inference, all on dissociated data",
    "Wednesday afternoon: spatial, both branches. Visium and Visium HD on the "
    "sequencing side, MERSCOPE, CosMx and CODEX on the imaging side.",
    "See wed_morning_tutorials.md and wed_afternoon_tutorials.md for the "
    "tutorial list behind each of these",
], size=20)

# --------------------------------------------------------------- questions

divider("Questions")

content("What cell types are here?", [
    "The discovery question: you do not yet know what is in the sample",
    "Dimension reduction, then graph-based clustering, then marker genes per "
    "cluster",
    "The resolution parameter is a choice, not a result. Two clusters or five "
    "is often defensible either way.",
    "Clusters are hypotheses about cell types, not cell types",
])

content("Which known cell type is each cluster?", [
    "The annotation question, and it is a different problem from finding the "
    "clusters",
    "By hand: score canonical markers per cluster and name it",
    "Automatically: classify against a reference atlas, or transfer labels "
    "from an annotated dataset",
    "Automated labels are confident even when they are wrong. Always check the "
    "markers you would have used by hand.",
], size=20)

content("Are the samples even comparable?", [
    "Usually has to be settled before clustering and annotation can be trusted",
    "Batch correction and integration: make cells of the same type from "
    "different samples, donors or chemistries land together",
    "Query-to-reference mapping: place new cells into an existing atlas rather "
    "than re-clustering from scratch",
    "The failure modes run both ways. Under-correct and every cluster is one "
    "sample, over-correct and you have erased the biological difference you "
    "came to find.",
], size=20)

content("How does one cell type differ between conditions, or over time?", [
    "The commonest question anyone actually asks of the data",
    "Subset to the cell type, then test condition against condition",
    "Testing cell by cell treats cells as independent replicates. They are "
    "not: cells from one donor are correlated, and p-values come out absurdly "
    "small.",
    "Pseudobulk instead: sum counts per cell type per sample, then use bulk "
    "RNA-seq machinery with samples as the unit of replication",
], size=20)

content("Are there more or fewer of a cell type?", [
    "Differential abundance, or compositional analysis",
    "Harder than it looks: proportions must sum to one, so one population "
    "expanding makes every other one appear to shrink",
    "Also depends entirely on how you drew the cluster boundaries",
    "Milo sidesteps the clustering by testing abundance in overlapping "
    "neighbourhoods of the k-nearest-neighbour graph instead",
], size=20)

content("How do cells order along a continuous process?", [
    "For differentiation, activation, or any gradual change, discrete clusters "
    "are the wrong model",
    "Trajectory inference orders cells along a path, and pseudotime says how "
    "far along a cell sits",
    "Pseudotime is a distance, not a clock. Units are arbitrary and direction "
    "comes from biology, not from the algorithm.",
    "RNA velocity uses the unspliced-to-spliced ratio to point at where a cell "
    "is heading, and is more fragile than its pictures suggest",
    "Wednesday morning is entirely about this",
], size=20)

content("What is driving the difference?", [
    "Gene set and pathway enrichment: is the gene list you found enriched for "
    "something already named",
    "Gene regulatory network inference: which transcription factors have their "
    "targets moving together, for example SCENIC",
    "In tumour work, inferred copy number from expression separates malignant "
    "cells from the normal cells beside them",
], size=20)

content("Which cells are talking to which?", [
    "Ligand-receptor analysis: a ligand expressed in one cell type and its "
    "receptor in another",
    "This is a prediction from co-expression. Nothing about contact, "
    "secretion, or an actual signalling event has been observed.",
    "Two cell types can express a matched pair and sit at opposite ends of the "
    "tissue",
    "Which is exactly why spatial data makes this question much stronger",
], size=20)

content("Where are the cells in the image?", [
    "Imaging-based spatial gives molecule coordinates, not cells. You have to "
    "draw the boundaries.",
    "Segmentation from a nuclear stain, expanded outward, or learned from cell "
    "membranes",
    "Every downstream count depends on it, and mis-assigned transcripts look "
    "exactly like doublets or like hybrid cell types that do not exist",
    "Spot-based data has the mirror-image problem: no segmentation to do, but "
    "each spot is a mixture, so you deconvolve it into cell-type proportions "
    "instead",
], size=20)

content("How is the tissue organised?", [
    "The question you can only ask once position is retained",
    "Spatial domains: recurring neighbourhoods of cell types, tumour core "
    "versus margin, germinal centre versus mantle",
    "Spatially variable genes: expression structured by location rather than "
    "by cell type",
    "Neighbourhood analysis: which cell types are found next to which, more "
    "often than chance",
], size=20)

content("Can I combine different measurements from the same sample?", [
    "Paired: RNA and protein, or RNA and ATAC, measured in the same cell. "
    "Alignment is given, the work is in weighting the modalities.",
    "Unpaired: the same tissue assayed two ways in different cells. You have "
    "to infer which cell corresponds to which.",
    "Mosaic: a set of datasets that overlap partially in modality, which is "
    "what atlas building really looks like",
    "Also the bridge between the two spatial branches, and between spatial and "
    "dissociated reference data",
], size=20)

content("The honest caveats:", [
    "Most of these outputs are inference, not measurement",
    "Cluster resolution is a choice you made, and the number of cell types "
    "moves with it",
    "Cell-cell interaction is a prediction from co-expression",
    "Pseudotime is a distance, not a clock",
    "Segmentation error propagates into every count downstream",
    "None of this is a reason not to do it. It is a reason to say what you did.",
], size=20)

content("References: analysis and best practice", [
    "Luecken and Theis 2019, Current best practices in single-cell RNA-seq "
    "analysis: a tutorial, Molecular Systems Biology 15(6):e8746. "
    "https://doi.org/10.15252/msb.20188746",
    "Heumos et al. 2023, Best practices for single-cell analysis across "
    "modalities, Nature Reviews Genetics 24(8):550-572. "
    "https://doi.org/10.1038/s41576-023-00586-w",
    "Zappia and Theis 2021, Over 1000 tools reveal trends in the single-cell "
    "RNA-seq analysis landscape, Genome Biology 22:301. "
    "https://doi.org/10.1186/s13059-021-02519-4",
], size=16)

content("References: technologies", [
    "Moses and Pachter 2022, Museum of spatial transcriptomics, Nature Methods "
    "19(5):534-546. https://doi.org/10.1038/s41592-022-01409-2",
    "Palla et al. 2022, Spatial components of molecular tissue biology, Nature "
    "Biotechnology 40(3):308-318. "
    "https://doi.org/10.1038/s41587-021-01182-1",
    "Dixit et al. 2016, Perturb-Seq: Dissecting Molecular Circuits with "
    "Scalable Single-Cell RNA Profiling of Pooled Genetic Screens, Cell "
    "167(7):1853-1866.e17. https://doi.org/10.1016/j.cell.2016.11.038",
    "Replogle et al. 2022, Mapping information-rich genotype-phenotype "
    "landscapes with genome-scale Perturb-seq, Cell 185(14):2559-2575.e28. "
    "https://doi.org/10.1016/j.cell.2022.05.013",
], size=16)

prs.save(OUT)
print("wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
